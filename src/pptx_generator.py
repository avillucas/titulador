import os
import re
import pptx
from pptx.util import Pt
from typing import Dict, Any, List, Tuple
from src.models import TituloData

DEFAULT_DASH_LINE = "--------------------------------------------------------------------------"

# Standard shape mapping based on template inspection
SHAPE_MAP_SLIDE1 = {
    85: "apellido_nombre",
    86: "documento",
    87: "horas_cursada",
    88: "titulo_linea1",
    89: "titulo_linea2",
    90: "resolucion",
    91: "emision_dia",
    92: "emision_mes",
    93: "emision_ano",
}

SHAPE_MAP_SLIDE2 = {
    99: 0,   # Slot 1  (Columna izquierda, fila 1)
    101: 1,  # Slot 2  (Columna izquierda, fila 2)
    103: 2,  # Slot 3  (Columna izquierda, fila 3)
    105: 3,  # Slot 4  (Columna izquierda, fila 4)
    107: 4,  # Slot 5  (Columna izquierda, fila 5)
    100: 5,  # Slot 6  (Columna derecha, fila 1)
    102: 6,  # Slot 7  (Columna derecha, fila 2)
    104: 7,  # Slot 8  (Columna derecha, fila 3)
    106: 8,  # Slot 9  (Columna derecha, fila 4)
    108: 9,  # Slot 10 (Columna derecha, fila 5)
    109: "fecha_egreso",
    110: "numero_egresado",
    111: "numero_cpf",
    112: "distrito_cpf",
}

def split_certificacion_title(title: str, max_line1: int = 38) -> Tuple[str, str]:
    """Splits a certification title so part 1 fits on Shape 89 and part 2 wraps to Shape 90."""
    title = title.strip()
    if not title or len(title) <= max_line1:
        return title, ""

    words = title.split()
    l1_words = []
    curr_len = 0
    for w in words:
        if curr_len + len(w) + (1 if l1_words else 0) <= max_line1:
            l1_words.append(w)
            curr_len += len(w) + 1
        else:
            break

    if l1_words and len(" ".join(l1_words)) >= max_line1 - 10:
        part1 = " ".join(l1_words)
        part2 = title[len(part1):].strip()
        return part1, part2

    # Syllable / Hyphen break fallback near max_line1
    cut = max_line1
    part1 = title[:cut] + "-"
    part2 = title[cut:].lstrip()
    return part1, part2


def format_module_text(module_str: str, max_single_line_chars: int = 52, max_line_chars: int = 42) -> str:
    """Formats module text with dot-testing per Circular 04-2020 f) and g)."""
    module_str = module_str.strip()
    if not module_str or module_str.startswith("----"):
        return DEFAULT_DASH_LINE

    # Check if single line fits completely (denom + code)
    if len(module_str) <= max_single_line_chars:
        num_dots = max(3, 60 - len(module_str))
        return f"{module_str} {'.' * num_dots}"

    # Extract module code if present (e.g. (MM0011.1))
    match = re.search(r'^(.*?)\s*(\([A-Z0-9\.]+\))?$', module_str)
    if match and match.group(2):
        denom, code = match.group(1).strip(), match.group(2).strip()
    else:
        denom = module_str
        code = ""

    # If denom fits on line 1 without code
    if len(denom) <= max_line_chars:
        num_dots = max(3, 60 - len(denom))
        line1 = f"{denom} {'.' * num_dots}"
        line2 = f" {code}" if code else ""
        return f"{line1}\n{line2}"

    # Denom itself spans 2 lines
    words = denom.split()
    l1_words = []
    curr = 0
    for w in words:
        if curr + len(w) + (1 if l1_words else 0) <= max_line_chars:
            l1_words.append(w)
            curr += len(w) + 1
        else:
            break

    if l1_words:
        part1 = " ".join(l1_words)
        part2 = denom[len(part1):].strip()
    else:
        part1 = denom[:max_line_chars]
        part2 = denom[max_line_chars:].strip()

    num_dots = max(3, 60 - len(part1))
    line1 = f"{part1} {'.' * num_dots}"
    line2 = f" {part2} {code}".strip()
    return f"{line1}\n {line2}"


def set_text_frame_content(shape, text: str, font_name: str = "Arial", font_size_pt: float = 12.0, bold: bool = False, word_wrap: bool = False):
    """Utility to reliably set text frame content, font properties, margins, and line spacing."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)

    lines = text.split("\n")
    tf.clear()

    for idx, line in enumerate(lines):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.font.name = font_name
        p.font.size = Pt(font_size_pt)
        p.font.bold = bold
        p.text = line


class PPTXGenerator:
    def __init__(self, template_path: str):
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template file not found at {template_path}")
        self.template_path = template_path

    def generate(self, data: TituloData, output_path: str) -> str:
        prs = pptx.Presentation(self.template_path)
        
        # Process Slide 1 (Frente)
        if len(prs.slides) > 0:
            slide1 = prs.slides[0]
            
            # Prepare Certificación title line 1 & line 2 split
            cert_part1, cert_part2 = split_certificacion_title(data.titulo_linea2, max_line1=32)
            
            for shape in slide1.shapes:
                if shape.shape_id == 85:  # Apellido y Nombre (Arial 14 Bold UPPERCASE)
                    val = str(data.apellido_nombre).upper()
                    sz = 12.0 if len(val) > 32 else 14.0
                    set_text_frame_content(shape, val, font_name="Arial", font_size_pt=sz, bold=True, word_wrap=False)
                elif shape.shape_id == 86:  # D.U. Documento (Arial 12)
                    set_text_frame_content(shape, str(data.documento), font_name="Arial", font_size_pt=12.0, word_wrap=False)
                elif shape.shape_id == 87:  # Horas (Arial 12) - Move slightly up to avoid line overlap
                    from pptx.util import Mm
                    shape.left = Mm(149.5)
                    shape.top = Mm(78.6)
                    set_text_frame_content(shape, str(data.horas_cursada), font_name="Arial", font_size_pt=12.0, word_wrap=False)
                elif shape.shape_id == 88:  # Trayecto / Curso (Arial 11-12)
                    from pptx.util import Mm
                    shape.top = Mm(79.0)
                    val = str(data.titulo_linea1)
                    sz = 10.5 if len(val) > 42 else 11.5
                    set_text_frame_content(shape, val, font_name="Arial", font_size_pt=sz, word_wrap=False)
                elif shape.shape_id == 89:  # Certificación DE Line 1 - Starts at 114mm (after CERTIFICADO DE) and ends at 202mm
                    from pptx.util import Mm
                    shape.left = Mm(114.0)
                    shape.top = Mm(87.5)
                    shape.width = Mm(88.0)  # Ends at 202mm
                    set_text_frame_content(shape, cert_part1, font_name="Arial", font_size_pt=12.0, word_wrap=False)



                elif shape.shape_id == 90:  # Certificación DE Line 2 + Resolución Nro.
                    res_text = str(data.resolucion)
                    if cert_part2:
                        full_line2 = f"{cert_part2} {res_text}"
                    else:
                        full_line2 = res_text
                    set_text_frame_content(shape, full_line2, font_name="Arial", font_size_pt=12.0, word_wrap=False)
                elif shape.shape_id == 91:  # Emisión Día
                    set_text_frame_content(shape, str(data.emision_dia), font_name="Arial", font_size_pt=12.0, word_wrap=False)
                elif shape.shape_id == 92:  # Emisión Mes
                    set_text_frame_content(shape, str(data.emision_mes), font_name="Arial", font_size_pt=12.0, word_wrap=False)
                elif shape.shape_id == 93:  # Emisión Año
                    set_text_frame_content(shape, str(data.emision_ano), font_name="Arial", font_size_pt=12.0, word_wrap=False)


        # Process Slide 2 (Dorso)
        if len(prs.slides) > 1:
            slide2 = prs.slides[1]
            for shape in slide2.shapes:
                if shape.shape_id in SHAPE_MAP_SLIDE2:
                    target = SHAPE_MAP_SLIDE2[shape.shape_id]
                    if isinstance(target, int):
                        # Module slot 0..9
                        if target < len(data.modulos) and data.modulos[target].strip():
                            formatted_mod = format_module_text(data.modulos[target].strip(), max_line_chars=44)
                            set_text_frame_content(shape, formatted_mod, font_name="Arial", font_size_pt=10.0, word_wrap=False)
                        else:
                            set_text_frame_content(shape, DEFAULT_DASH_LINE, font_name="Arial", font_size_pt=10.0, word_wrap=False)
                    elif target == "fecha_egreso":
                        set_text_frame_content(shape, str(data.fecha_egreso), font_name="Arial", font_size_pt=12.0, word_wrap=False)
                    elif target == "numero_egresado":
                        set_text_frame_content(shape, str(data.numero_egresado), font_name="Arial", font_size_pt=12.0, word_wrap=False)
                    elif target == "numero_cpf":
                        set_text_frame_content(shape, str(data.numero_cpf), font_name="Arial", font_size_pt=12.0, word_wrap=False)
                    elif target == "distrito_cpf":
                        set_text_frame_content(shape, str(data.distrito_cpf), font_name="Arial", font_size_pt=12.0, word_wrap=False)

        # Ensure output directory exists
        os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
        prs.save(output_path)
        return output_path

