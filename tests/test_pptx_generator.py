import unittest
import os
import pptx
from src.models import TituloData
from src.pptx_generator import PPTXGenerator

class TestPPTXGenerator(unittest.TestCase):
    def test_dash_padding_in_shapes(self):
        data = TituloData(
            apellido_nombre='VILLALBA JORGE HUMBERTO',
            documento='17.924.873',
            horas_cursada='540',
            titulo_linea1='Electricista en Inmuebles',
            titulo_linea2='Electricista en Inmuebles',
            resolucion='Resolución Nro. RESFC-2019-6908-GDEBA-DGCYE / RESOL-2022-2139-APN-ME--',
            emision_dia='09',
            emision_mes='Julio',
            emision_ano='26',
            modulos=['Módulo 1', 'Módulo 2'],
            fecha_egreso='09 de Junio de 2026',
            numero_egresado='1',
            numero_cpf='412',
            distrito_cpf='Lomas de Zamora'
        )

        gen = PPTXGenerator()
        out_path = 'output/test_dash_padding.pptx'
        gen.generate(data, out_path)

        self.assertTrue(os.path.exists(out_path))

        prs = pptx.Presentation(out_path)
        slide1 = prs.slides[0]

        shape_dict = {shape.shape_id: shape.text_frame.text for shape in slide1.shapes if shape.has_text_frame}

        # Shape 88: Title line 1 before hours should be padded with hyphens
        self.assertIn('Electricista en Inmuebles-', shape_dict[88])
        self.assertTrue(shape_dict[88].endswith('-'))

        # Shape 89: Certificacion line 1 should be padded with hyphens
        self.assertIn('Electricista en Inmuebles-', shape_dict[89])
        self.assertTrue(shape_dict[89].endswith('-'))

        # Shape 90: Resolution line should end with hyphens
        self.assertTrue(shape_dict[90].endswith('--'))

        # Clean up test output
        if os.path.exists(out_path):
            os.remove(out_path)

if __name__ == '__main__':
    unittest.main()
